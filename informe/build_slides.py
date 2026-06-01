# -*- coding: utf-8 -*-
"""Genera informe/slides.pptx para la presentacion del TP (Inteligencia Colectiva)."""
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# Requiere: pip install python-pptx Pillow. Regenerar:  python informe/build_slides.py
INF = pathlib.Path(__file__).resolve().parent
ROOT = INF.parent
OUT = INF / "slides.pptx"

IMG_EXTRACT = INF / "Capa de Extracción y-2026-05-31-155807.png"
IMG_DER = INF / "DER.png"
IMG_ETL = INF / "ETL.png"

NAVY   = RGBColor(0x0F, 0x2A, 0x43)
TEAL   = RGBColor(0x2A, 0x9D, 0x8F)
GRAY   = RGBColor(0x6B, 0x72, 0x80)
CODEBG = RGBColor(0xF3, 0xF4, 0xF6)
CODEFG = RGBColor(0x11, 0x18, 0x27)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHTLINE = RGBColor(0xD1, 0xD5, 0xDB)
FONT = "Calibri"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5


def _tb(slide, l, t, w, h):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    return box, box.text_frame


def _set(p, text, size, color, bold=False, font=FONT, align=PP_ALIGN.LEFT):
    p.text = text
    p.alignment = align
    r = p.runs[0]
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = font
    r.font.color.rgb = color
    return p


def new_slide(kicker=None, title=None, page=None):
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.16), Inches(SH))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
    if kicker is not None:
        _, tf = _tb(s, 0.6, 0.34, 12.2, 0.4)
        _set(tf.paragraphs[0], kicker.upper(), 13, TEAL, bold=True)
    if title is not None:
        _, tf = _tb(s, 0.58, 0.66, 12.2, 1.05)
        _set(tf.paragraphs[0], title, 30, NAVY, bold=True)
    _, ftf = _tb(s, 0.6, 7.02, 9.5, 0.35)
    _set(ftf.paragraphs[0], "Inteligencia Colectiva. TP Integrador BDD. Grupo 10", 10, GRAY)
    if page is not None:
        _, ptf = _tb(s, 12.2, 7.02, 0.9, 0.35)
        _set(ptf.paragraphs[0], str(page), 10, GRAY, align=PP_ALIGN.RIGHT)
    return s


def bullets(slide, items, l, t, w, h, size=18, gap=7):
    _, tf = _tb(slide, l, t, w, h)
    for i, it in enumerate(items):
        text, lvl = it if isinstance(it, tuple) else (it, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = "•  " if lvl == 0 else "–  "
        _set(p, prefix + text, size if lvl == 0 else size - 2,
             NAVY if lvl == 0 else GRAY)
        p.space_after = Pt(gap); p.level = lvl
    return tf


def code_box(slide, code, l, t, w, h, size=13):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = CODEBG
    shp.line.color.rgb = LIGHTLINE; shp.line.width = Pt(0.75); shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.12); tf.margin_bottom = Inches(0.12)
    for i, ln in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set(p, ln if ln else " ", size, CODEFG, font=MONO); p.space_after = Pt(1)
    return shp


def image_fit(slide, path, l, t, maxw, maxh):
    if not pathlib.Path(path).exists():
        _, tf = _tb(slide, l, t, maxw, maxh)
        _set(tf.paragraphs[0], "[falta imagen: %s]" % pathlib.Path(path).name, 12, GRAY)
        return
    iw, ih = Image.open(path).size
    ar = iw / ih; boxar = maxw / maxh
    if ar > boxar:
        w = maxw; h = maxw / ar
    else:
        h = maxh; w = maxh * ar
    pic = slide.shapes.add_picture(str(path), Inches(l + (maxw - w) / 2), Inches(t + (maxh - h) / 2),
                                   Inches(w), Inches(h))
    pic.line.color.rgb = LIGHTLINE; pic.line.width = Pt(1)
    return pic


def table(slide, headers, rows, l, t, w, colw=None, size=14):
    nr, nc = len(rows) + 1, len(headers)
    gt = slide.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(0.3 * nr)).table
    if colw:
        for j, cw in enumerate(colw):
            gt.columns[j].width = Inches(cw)
    for j, htxt in enumerate(headers):
        c = gt.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = TEAL
        _set(c.text_frame.paragraphs[0], htxt, size, WHITE, bold=True)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else RGBColor(0xF7, 0xF9, 0xFA)
            _set(c.text_frame.paragraphs[0], str(val), size - 1, NAVY)
    return gt


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


page = 0

# --- Portada ---
s = prs.slides.add_slide(BLANK)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SW), Inches(SH))
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.55), Inches(SW), Inches(0.07))
acc.fill.solid(); acc.fill.fore_color.rgb = TEAL; acc.line.fill.background()
_, tf = _tb(s, 0.9, 2.0, 11.5, 1.4)
_set(tf.paragraphs[0], "Inteligencia Colectiva", 54, WHITE, bold=True)
_, tf = _tb(s, 0.92, 3.35, 11.5, 1.0)
_set(tf.paragraphs[0], "Data Warehouse políglota: MongoDB + PostgreSQL", 26, TEAL, bold=True)
_, tf = _tb(s, 0.92, 4.8, 11.6, 1.7)
for i, line in enumerate([
        "TP Integrador de Bases de Datos. Grupo 10.",
        "Fraga · Piuselli · Arnesano · Kronenberger · Korenblit",
        "Corre en la nube: Supabase (PostgreSQL) y MongoDB Atlas."]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    _set(p, line, 16, WHITE if i != 2 else RGBColor(0xB9, 0xE3, 0xDD)); p.space_after = Pt(7)
notes(s, "Modelo de chofer único: uno comparte pantalla, el resto relata. 6 min por persona, cronómetro a la vista. "
         "Orden: Damián, Matías, Gabi, Julián, Tomás.")

# --- S1 Escenario ---
page += 1
s = new_slide("Escenario", "La empresa y el problema", page)
bullets(s, [
    "Inteligencia Colectiva: consultora de opinión pública.",
    "Encuesta por teléfono (IVR) y por formularios web.",
    "El cuello de botella es la ingesta: muchas encuestas, de varios canales, todas distintas.",
    ("Cada cuestionario trae otra cantidad y otro tipo de preguntas.", 1),
    ("El formato cambia seguido, así que no hay un esquema fijo.", 1),
    "El cliente quiere tableros de tendencias, cortes por edad, NSE y región, y proyecciones.",
], 0.7, 1.95, 11.9, 4.8, size=19)
notes(s, "(Damián) Plantear el dolor del negocio: datos crudos, multicanal, esquemas que cambian. Eso obliga a dos "
         "tratamientos distintos: guardar flexible y analizar con rigor.")

# --- S2 Solución ---
page += 1
s = new_slide("Escenario", "La solución: arquitectura políglota", page)
bullets(s, [
    "El flujo tiene dos momentos muy distintos.",
    ("Primero entra todo crudo y desordenado.", 1),
    ("Después se ordena para analizar y predecir.", 1),
    "Usamos un motor para cada momento.",
    ("MongoDB guarda lo crudo (Data Lake).", 1),
    ("PostgreSQL hace el análisis (Data Warehouse).", 1),
    "El recorrido es extracción, transformación y presentación.",
], 0.7, 1.95, 6.5, 4.8, size=18)
image_fit(s, IMG_EXTRACT, 7.5, 1.95, 5.3, 4.6)
notes(s, "(Damián) La decisión central del TP. Mostrar el diagrama de las tres capas y dejar la frase: usar el motor "
         "correcto para el trabajo correcto.")

# --- S3 MongoDB ---
page += 1
s = new_slide("Motores", "MongoDB: el Data Lake (NoSQL)", page)
bullets(s, [
    "Proveedor: MongoDB Inc. Licencia SSPL (libre; Atlas es la versión gestionada).",
    "Lo elegimos porque es schema-less: absorbe encuestas que cambian de forma sin migrar nada.",
    "El documento JSON/BSON encaja con la jerarquía pregunta-respuesta.",
    "Se consigue local, con Docker (docker pull mongo), o en la nube con Atlas.",
], 0.7, 1.95, 11.9, 4.8, size=18)
notes(s, "(Matías) Concepto a remarcar: esquema flexible. La estructura real del documento va en la slide que sigue.")

# --- Documento Mongo ---
page += 1
s = new_slide("Motores", "El documento en MongoDB", page)
bullets(s, [
    "Dos colecciones: surveys (la encuesta y sus preguntas) y responses (cada respuesta).",
    "Una respuesta llega anidada, tal cual la guarda Mongo:",
], 0.7, 1.85, 12.0, 1.3, size=17)
code_box(s,
    '{\n'
    '  "_id": "resp_0001_001",\n'
    '  "encuesta_id": "survey_0001",\n'
    '  "encuestado_id": "person_00017",\n'
    '  "fecha": "2025-03-14T10:32:00",\n'
    '  "fuente": "IVR",\n'
    '  "respuestas": [\n'
    '    { "pregunta_id": "survey_0001_q01",\n'
    '      "opcion_id": "survey_0001_q01_opt03",\n'
    '      "valor": 3, "texto": "Movimiento Popular" },\n'
    '    { "pregunta_id": "survey_0001_q02", "valor": 2, "texto": "Buena" }\n'
    '  ]\n'
    '}',
    0.7, 3.15, 11.9, 3.5, size=13)
notes(s, "(Matías) Estructura real de la coleccion responses: encuesta_id, encuestado_id, fecha, fuente y el arreglo "
         "respuestas con cada pregunta. El documento anidado es la razon por la que hace falta un esquema flexible.")

# --- S4 PostgreSQL ---
page += 1
s = new_slide("Motores", "PostgreSQL: el Data Warehouse (SQL)", page)
bullets(s, [
    "Proveedor: PostgreSQL Global Development Group. Licencia tipo BSD/MIT, gratis.",
    "Lo elegimos por el rigor: ACID, y funciones analíticas que corren dentro del motor.",
    "La segmentación y la predicción bayesiana pasan en SQL, no en la app.",
    "Lo usamos vía Supabase, que lo deja andando en la nube.",
    "Se consigue local, con Docker, o gestionado (Supabase, RDS, Cloud SQL).",
], 0.7, 1.95, 11.9, 4.8, size=18)
notes(s, "(Gabi) Justificar el lado SQL: por qué un motor relacional para la capa analítica. Delegamos la matemática "
         "pesada al motor, no a la aplicación.")

# --- S5 Un motor para cada cosa ---
page += 1
s = new_slide("Motores", "Un motor para cada cosa", page)
bullets(s, [
    "MongoDB guarda lo crudo. PostgreSQL lo analiza.",
    "La regla que seguimos: usar el motor correcto para el trabajo correcto.",
    "Los dos están gestionados en la nube: Postgres en Supabase, Mongo en Atlas.",
    ("Eso nos pone en el tramo más alto de la consigna (nube).", 1),
    "Las credenciales van en .env, fuera del código. En Supabase sumamos roles y RLS.",
], 0.7, 1.95, 11.9, 4.8, size=19)
notes(s, "(Damián) Cierre del bloque de motores y del overview de stack. Dejar claro que está todo en la nube, que es "
         "lo que pide la nota máxima.")

# --- S6 Modelado ---
page += 1
s = new_slide("Modelado", "Estrella, con una rama en copo de nieve", page)
bullets(s, [
    "Hechos: fact_survey_responses, una fila por cada respuesta.",
    "Dimensiones: surveys, questions, answer_options, respondents y time.",
    "No es estrella pura: answer_options cuelga de questions, y questions de surveys.",
    ("Esa rama normalizada es lo que la vuelve híbrida (copo de nieve).", 1),
    "Auditoría aparte: etl_process_executions guarda cada corrida del ETL.",
], 0.7, 1.95, 6.3, 4.8, size=18)
image_fit(s, IMG_DER, 7.3, 1.85, 5.5, 4.9)
notes(s, "(Gabi) Mostrar el DER. Justificar el híbrido: estrella para agregar rápido, pero answer_options -> questions "
         "-> surveys está normalizado (copo de nieve). Nombrar la tabla de auditoría.")

# --- S7 Hash ---
page += 1
s = new_slide("Modelado", "Perfilado determinístico con SHA-256", page)
bullets(s, [
    "El perfil del encuestado no se guarda: se calcula del ID con un hash.",
    "La misma persona da siempre el mismo perfil, sin almacenar datos personales.",
    "Así tenemos región, NSE y edad consistentes entre encuestas.",
], 0.7, 1.9, 12.0, 1.7, size=18)
code_box(s,
    'def profile_for(respondent_id: str) -> dict:\n'
    '    h = hashlib.sha256(respondent_id.encode()).digest()\n'
    '    region = REGIONS[h[0] % len(REGIONS)]\n'
    '    nse    = NSE_LEVELS[h[1] % len(NSE_LEVELS)]\n'
    '    edad   = 16 + (int.from_bytes(h[3:5], "big") % 75)\n'
    '    return {"region": region, "nse": nse, "edad": edad, ...}',
    0.7, 3.7, 11.9, 2.6, size=14)
_, tf = _tb(s, 0.72, 6.32, 11.9, 0.4)
_set(tf.paragraphs[0], "src/demographics.py", 12, GRAY, font=MONO)
notes(s, "(Gabi) El punto para lucirse. El hash del respondent_id deriva región, NSE, edad y género. No guardamos "
         "datos personales: los reconstruimos. Es el snippet real del repo.")

# --- S8 ETL flujo ---
page += 1
s = new_slide("ETL", "El proceso ETL", page)
bullets(s, [
    "El ETL lee de Mongo por lotes, aplana el JSON, lo enriquece con el hash y carga en Postgres.",
    "Está en src/etl_mongo_to_postgres.py y corre incremental: no borra lo que ya está.",
    "Cada corrida queda registrada en etl_process_executions.",
], 0.7, 1.95, 5.9, 4.6, size=18)
image_fit(s, IMG_ETL, 6.8, 1.85, 6.0, 4.9)
notes(s, "(Julián) Mostrar el diagrama de flujo y después pasar al IDE. Explicar la extracción y el aplanado del JSON "
         "anidado, más el enriquecido por hash.")

# --- S9 Idempotencia ---
page += 1
s = new_slide("ETL", "Idempotencia: ON CONFLICT", page)
bullets(s, [
    "El ETL se puede correr mil veces y no duplica nada.",
    "En las dimensiones hace upsert: si ya existe, lo actualiza (DO UPDATE).",
    "En los hechos ignora el duplicado (DO NOTHING), apoyado en la unicidad uq_fact_response_question.",
], 0.7, 1.9, 12.0, 1.7, size=18)
code_box(s,
    '# Dimensiones: upsert (crea o actualiza)\n'
    'stmt.on_conflict_do_update(\n'
    '    index_elements=["respondent_id"],\n'
    '    set_={"region": stmt.excluded.region, ...})\n'
    '\n'
    '# Hechos: ignora duplicados\n'
    'stmt.on_conflict_do_nothing(index_elements=["date_key"])',
    0.7, 3.6, 11.9, 2.7, size=14)
notes(s, "(Julián) Idempotente quiere decir que correrlo de nuevo da el mismo estado. DO UPDATE actualiza (por ej. el "
         "status active a closed), DO NOTHING evita duplicar en la fact.")

# --- S10 Busquedas ---
page += 1
s = new_slide("Consultas", "Búsquedas por una y por dos claves", page)
bullets(s, [
    "La consigna pide buscar por una clave y por dos.",
    "Salen de sql/busquedas.sql, ya sobre el Data Warehouse cargado.",
], 0.7, 1.9, 12.0, 1.3, size=18)
code_box(s,
    "-- una clave (PK survey_id)\n"
    "select * from dim_surveys\n"
    "where survey_id = 'survey_0001';\n"
    "\n"
    "-- dos claves (clave compuesta de la fact)\n"
    "select * from fact_survey_responses\n"
    "where response_id = 'resp_0001_001'\n"
    "  and question_id = 'survey_0001_q01';",
    0.7, 3.25, 11.9, 3.0, size=14)
notes(s, "(Gabi) Cubre el ítem 4.5 de la consigna. Una clave: la PK de dim_surveys. Dos claves: la fact tiene clave "
         "compuesta (response_id, question_id), uq_fact_response_question.")

# --- S11 Segmentacion ---
page += 1
s = new_slide("Minería", "Función de segmentación", page)
bullets(s, [
    "segmentar_respuestas(categoria, desde, hasta, dimension).",
    "Agrupa las respuestas por la dimensión que le pidas: región, NSE, y demás.",
    "Cuenta y saca el porcentaje de cada segmento, todo dentro de SQL.",
], 0.7, 1.9, 12.0, 1.7, size=17)
code_box(s,
    "select * from segmentar_respuestas(\n"
    "  'intencion_voto',\n"
    "  '2024-01-01', '2026-12-31', 'region');",
    0.7, 3.55, 6.2, 1.5, size=13)
table(s, ["segmento", "cantidad", "%"], [
    ["NEA", 689, "17.23"], ["AMBA", 683, "17.08"], ["Centro", 678, "16.95"],
    ["NOA", 666, "16.65"], ["Patagonia", 661, "16.53"], ["Cuyo", 623, "15.58"],
], 7.3, 3.05, 5.5, colw=[2.2, 1.7, 1.6], size=14)
notes(s, "(Tomás) Ítem 5.1 de la consigna. Para qué sirve: ver cómo opina cada región o NSE sin reescribir queries. "
         "Mostrar la función y su resultado real, suma cercana a 100%.")

# --- S12 Prediccion ---
page += 1
s = new_slide("Minería", "Función de predicción (bayesiana)", page)
bullets(s, [
    "predecir_shares(categoria, fecha_corte, lambda, prior, region, nse).",
    "Es un Dirichlet-Multinomial con descuento por recencia: lambda le da más peso a lo nuevo.",
    "Devuelve el share estimado y su intervalo de credibilidad del 95%.",
], 0.7, 1.9, 12.0, 1.7, size=17)
code_box(s,
    "select * from predecir_shares(\n"
    "  'intencion_voto', '2026-09-26', 0.0, 5.0);",
    0.7, 3.55, 6.2, 1.2, size=13)
table(s, ["partido", "share", "IC 95%"], [
    ["Frente Federal", "26.2%", "24.8 a 27.5"],
    ["Unión Republicana", "21.7%", "20.4 a 23.0"],
    ["Movimiento Popular", "20.8%", "19.6 a 22.1"],
    ["Alianza Verde", "16.9%", "15.8 a 18.1"],
    ["Partido Vecinal", "14.4%", "13.3 a 15.4"],
], 7.3, 3.05, 5.5, colw=[2.5, 1.4, 1.6], size=13)
notes(s, "(Tomás) Ítem 5.2. Lambda es el olvido por día. Con más datos frescos, la banda de error se angosta. Toda la "
         "estadística vive en el motor SQL.")

# --- S13 Dashboard ---
page += 1
s = new_slide("Dashboard", "El dashboard, en vivo", page)
bullets(s, [
    "dashboard/app.py se conecta directo al Data Warehouse. Tiene más de cuatro vistas:",
    ("Intención de voto: barras con intervalo (predecir_shares) y un mapa de provincias por partido.", 1),
    ("Imagen de gobierno mes a mes, cruzando la fact con dim_time.", 1),
    ("Composición de la muestra: región, edad, género y canal.", 1),
    ("Explorador: movés el lambda y los filtros, y el modelo se recalcula al toque.", 1),
    "En la demo subimos el lambda y se ve cómo se abre el intervalo.",
], 0.7, 1.95, 11.9, 4.8, size=17)
notes(s, "(Tomás) Pasar al navegador (localhost:8501). Mostrar la predicción, el mapa de provincias, la evolución "
         "temporal y el explorador. Mover el slider lambda en la tab Predicción para que se vea en vivo.")

# --- S14 Cierre ---
page += 1
s = new_slide("Cierre", "Cómo cierra con el problema del cliente", page)
bullets(s, [
    "Volvemos al principio: encuestas crudas de muchos canales, y un tablero que las vuelve útiles.",
    "Cada motor en su lugar: Mongo para guardar, Postgres para analizar y predecir.",
    "La minería corre en la base, el ETL es idempotente y el perfil no expone datos personales.",
    "Está todo gestionado en la nube.",
    "Gracias. Preguntas.",
], 0.7, 1.95, 11.9, 4.8, size=20)
notes(s, "(Tomás) Cerrar volviendo al dolor del Bloque 1: la arquitectura políglota resuelve la ingesta flexible y el "
         "análisis con rigor. Abrir a preguntas.")

# --- S15 Apendice ---
page += 1
s = new_slide("Apéndice", "Tamaño y seguridad", page)
bullets(s, [
    "Una encuesta grande deja cerca de 1.6M filas en la fact, unos 50 MB.",
    ("Con 100 GB de disco aguanta años sin necesitar sharding.", 1),
    "Las credenciales van en .env, fuera del código.",
    ("En Supabase usamos RLS y roles separados: lectura para el dashboard, escritura para el ETL.", 1),
], 0.7, 1.95, 11.9, 4.6, size=18)
notes(s, "Slide de respaldo, por si preguntan por escala o seguridad.")

prs.save(str(OUT))
print("OK ->", OUT)
